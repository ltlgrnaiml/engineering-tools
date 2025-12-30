"""Image renderer for inserting images into PowerPoint shapes.

Per ADR-0022: ImageRenderer handles 'image_' category shapes.
Per ADR-0029: Uses shared ImageSpec contract.
"""

import io
import logging
from pathlib import Path
from typing import Any
from uuid import uuid4

from pptx.util import Inches

from shared.contracts.core.rendering import ImageSpec

from apps.pptx_generator.backend.core.shape_name_parser import ParsedShapeNameV2
from apps.pptx_generator.backend.renderers.base import BaseRenderer, RenderContext

logger = logging.getLogger(__name__)


class ImageRenderer(BaseRenderer):
    """Renderer for image shapes that insert images into placeholders.

    Handles shapes with names like:
    - image_logo
    - image_wafermap
    - image_chart_export
    """

    def __init__(self) -> None:
        """Initialize the image renderer."""
        super().__init__()
        self._image_cache: dict[str, bytes] = {}

    def can_render(self, parsed_name: ParsedShapeNameV2) -> bool:
        """Check if this is an image shape.

        Args:
            parsed_name: Parsed shape name.

        Returns:
            True if renderer is 'image'.
        """
        return parsed_name.renderer == "image"

    async def render(self, context: RenderContext) -> None:
        """Render image into shape.

        Args:
            context: Rendering context.
        """
        shape = context.shape
        parsed_name = context.parsed_name

        self.logger.info(f"[IMAGE] Starting render for {shape.name}")

        # Get image source from shape options or data
        image_source = self._get_image_source(parsed_name, context)

        if not image_source:
            self.logger.warning(f"[IMAGE] No image source found for {shape.name}")
            return

        try:
            # Load image data
            image_data = await self._load_image(image_source, context)

            if not image_data:
                self.logger.warning(f"[IMAGE] Failed to load image for {shape.name}")
                return

            # Insert image into shape
            self._insert_image(shape, image_data, parsed_name)

            self.logger.info(f"[IMAGE] Successfully rendered image into {shape.name}")

        except Exception as e:
            self.logger.error(f"[IMAGE] Failed to render image: {e}", exc_info=True)

    def _get_image_source(
        self,
        parsed_name: ParsedShapeNameV2,
        context: RenderContext,
    ) -> str | None:
        """Get image source path or URL from shape configuration.

        Args:
            parsed_name: Parsed shape name.
            context: Rendering context.

        Returns:
            Image source path or URL, or None if not found.
        """
        # Check for explicit path in options
        if parsed_name.options and "path" in parsed_name.options:
            return parsed_name.options["path"]

        # Check for image reference in data field
        if parsed_name.data:
            data_ref = parsed_name.data
            if isinstance(data_ref, str):
                # Could be a column name containing image paths
                if data_ref in context.data.columns:
                    # Get first non-null value
                    values = context.data[data_ref].dropna()
                    if len(values) > 0:
                        return str(values.iloc[0])
                # Or it could be a direct path/identifier
                return data_ref

        # Check metadata for image mappings
        if context.metadata and "images" in context.metadata:
            image_id = parsed_name.data or parsed_name.renderer
            if image_id in context.metadata["images"]:
                return context.metadata["images"][image_id]

        return None

    async def _load_image(
        self,
        source: str,
        context: RenderContext,
    ) -> bytes | None:
        """Load image data from source.

        Args:
            source: Image path or URL.
            context: Rendering context.

        Returns:
            Image bytes, or None if loading failed.
        """
        # Check cache first
        if source in self._image_cache:
            return self._image_cache[source]

        image_data: bytes | None = None

        # Try loading from file path
        if self._is_file_path(source):
            image_data = self._load_from_file(source, context)

        # Try loading from URL
        elif source.startswith(("http://", "https://")):
            image_data = await self._load_from_url(source)

        # Try loading from base64
        elif source.startswith("data:image"):
            image_data = self._load_from_base64(source)

        # Cache successful loads
        if image_data:
            self._image_cache[source] = image_data

        return image_data

    def _is_file_path(self, source: str) -> bool:
        """Check if source looks like a file path."""
        return (
            source.startswith(("/", "\\", "./", ".\\"))
            or (len(source) > 1 and source[1] == ":")  # Windows drive letter
            or source.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg"))
        )

    def _load_from_file(
        self,
        path: str,
        context: RenderContext,
    ) -> bytes | None:
        """Load image from file path.

        Args:
            path: File path.
            context: Rendering context.

        Returns:
            Image bytes, or None if loading failed.
        """
        # Resolve relative paths
        file_path = Path(path)
        if not file_path.is_absolute() and context.output_dir:
            file_path = context.output_dir / path

        if not file_path.exists():
            self.logger.warning(f"[IMAGE] File not found: {file_path}")
            return None

        try:
            return file_path.read_bytes()
        except Exception as e:
            self.logger.error(f"[IMAGE] Failed to read file: {e}")
            return None

    async def _load_from_url(self, url: str) -> bytes | None:
        """Load image from URL.

        Args:
            url: Image URL.

        Returns:
            Image bytes, or None if loading failed.
        """
        try:
            import aiohttp

            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=30) as response:
                    if response.status == 200:
                        return await response.read()
                    self.logger.warning(
                        f"[IMAGE] Failed to fetch URL: {url} (status {response.status})"
                    )
        except ImportError:
            self.logger.warning("[IMAGE] aiohttp not available for URL loading")
        except Exception as e:
            self.logger.error(f"[IMAGE] Failed to fetch URL: {e}")

        return None

    def _load_from_base64(self, data_uri: str) -> bytes | None:
        """Load image from base64 data URI.

        Args:
            data_uri: Base64 data URI.

        Returns:
            Image bytes, or None if parsing failed.
        """
        try:
            import base64

            # Parse data URI: data:image/png;base64,<data>
            header, data = data_uri.split(",", 1)
            return base64.b64decode(data)
        except Exception as e:
            self.logger.error(f"[IMAGE] Failed to parse base64 data: {e}")
            return None

    def _insert_image(
        self,
        shape: Any,
        image_data: bytes,
        parsed_name: ParsedShapeNameV2,
    ) -> None:
        """Insert image into PowerPoint shape.

        Args:
            shape: PowerPoint shape.
            image_data: Image bytes.
            parsed_name: Parsed shape name.
        """
        # Get shape dimensions and position
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        shape_name = shape.name

        # Create image stream
        image_stream = io.BytesIO(image_data)

        # Get parent slide shapes collection
        slide_shapes = shape._parent

        # Remove old shape
        sp = shape.element
        sp.getparent().remove(sp)

        # Add image as picture, maintaining aspect ratio by default
        fit_mode = parsed_name.options.get("fit", "contain") if parsed_name.options else "contain"

        if fit_mode == "stretch":
            # Stretch to fill
            pic = slide_shapes.add_picture(image_stream, left, top, width=width, height=height)
        else:
            # Contain: fit within bounds maintaining aspect ratio
            pic = slide_shapes.add_picture(image_stream, left, top)

            # Calculate scale to fit
            scale_x = width / pic.width
            scale_y = height / pic.height
            scale = min(scale_x, scale_y)

            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)

            # Center within original bounds
            pic.left = left + (width - pic.width) // 2
            pic.top = top + (height - pic.height) // 2

        # Restore shape name
        pic.name = shape_name

        self.logger.debug(f"[IMAGE] Inserted image at ({left}, {top}) with size ({pic.width}, {pic.height})")

    def build_image_spec(
        self,
        source_path: str,
        name: str,
        fit_mode: str = "contain",
    ) -> ImageSpec:
        """Build an ImageSpec from configuration per ADR-0029.

        Args:
            source_path: Path to image file.
            name: Spec name/identifier.
            fit_mode: Fit mode (stretch, contain, cover, none).

        Returns:
            ImageSpec conforming to shared rendering contract.
        """
        return ImageSpec(
            spec_id=f"image_{uuid4().hex[:8]}",
            name=name,
            source_path=source_path,
            fit_mode=fit_mode,
            source_tool="pptx",
        )
