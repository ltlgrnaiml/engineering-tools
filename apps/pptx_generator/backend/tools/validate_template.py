"""Template validation tool for checking shape naming conventions."""

import sys
from pathlib import Path

from pptx import Presentation

from apps.pptx_generator.backend.core.shape_name_parser import ValidationError, validate_shape_name


class TemplateValidationReport:
    """Report of template validation results."""

    def __init__(self, template_path: Path):
        """
        Initialize validation report.

        Args:
            template_path: Path to the template file.
        """
        self.template_path = template_path
        self.valid_shapes: list[tuple[int, str, str]] = []
        self.invalid_shapes: list[tuple[int, str, str, str]] = []
        self.warnings: list[tuple[int, str, str, list[str]]] = []

    def add_valid(self, slide_idx: int, shape_name: str, shape_id: str) -> None:
        """Add a valid shape to the report."""
        self.valid_shapes.append((slide_idx, shape_name, shape_id))

    def add_invalid(self, slide_idx: int, shape_name: str, shape_id: str, error: str) -> None:
        """Add an invalid shape to the report."""
        self.invalid_shapes.append((slide_idx, shape_name, shape_id, error))

    def add_warnings(
        self, slide_idx: int, shape_name: str, shape_id: str, warnings: list[str]
    ) -> None:
        """Add warnings for a shape to the report."""
        if warnings:
            self.warnings.append((slide_idx, shape_name, shape_id, warnings))

    def is_valid(self) -> bool:
        """Check if template has no errors (warnings are OK)."""
        return len(self.invalid_shapes) == 0

    def print_report(self) -> None:
        """Print formatted validation report."""
        print(f"\n{'=' * 80}")
        print(f"Template Validation Report: {self.template_path.name}")
        print(f"{'=' * 80}\n")

        print(f"✅ Valid shapes: {len(self.valid_shapes)}")
        print(f"⚠️  Shapes with warnings: {len(self.warnings)}")
        print(f"❌ Invalid shapes: {len(self.invalid_shapes)}\n")

        if self.valid_shapes:
            print(f"\n{'─' * 80}")
            print("✅ VALID SHAPES")
            print(f"{'─' * 80}")
            for slide_idx, shape_name, shape_id in self.valid_shapes:
                print(f"  Slide {slide_idx + 1}, Shape '{shape_id}': {shape_name}")

        if self.warnings:
            print(f"\n{'─' * 80}")
            print("⚠️  WARNINGS")
            print(f"{'─' * 80}")
            for slide_idx, shape_name, shape_id, warning_list in self.warnings:
                print(f"  Slide {slide_idx + 1}, Shape '{shape_id}': {shape_name}")
                for warning in warning_list:
                    print(f"    ⚠️  {warning}")

        if self.invalid_shapes:
            print(f"\n{'─' * 80}")
            print("❌ ERRORS")
            print(f"{'─' * 80}")
            for slide_idx, shape_name, shape_id, error in self.invalid_shapes:
                print(f"  Slide {slide_idx + 1}, Shape '{shape_id}': {shape_name}")
                print(f"    ❌ {error}")

        print(f"\n{'=' * 80}")
        if self.is_valid():
            print("✅ Template validation PASSED (warnings are informational)")
        else:
            print("❌ Template validation FAILED")
        print(f"{'=' * 80}\n")


def validate_template(template_path: Path) -> TemplateValidationReport:
    """
    Validate all shape names in a PowerPoint template.

    Args:
        template_path: Path to the PowerPoint template file.

    Returns:
        TemplateValidationReport with validation results.

    Raises:
        FileNotFoundError: If template file doesn't exist.
        ValueError: If template file is invalid.
    """
    if not template_path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")

    try:
        prs = Presentation(str(template_path))
    except Exception as e:
        raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

    report = TemplateValidationReport(template_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            shape_name = shape.name if hasattr(shape, "name") else f"Shape_{shape.shape_id}"
            shape_id = str(shape.shape_id)

            if not should_validate_shape(shape_name):
                continue

            try:
                parsed, warnings = validate_shape_name(shape_name)
                report.add_valid(slide_idx, shape_name, shape_id)
                if warnings:
                    report.add_warnings(slide_idx, shape_name, shape_id, warnings)
            except ValidationError as e:
                report.add_invalid(slide_idx, shape_name, shape_id, str(e))

    return report


def should_validate_shape(shape_name: str) -> bool:
    """
    Determine if a shape name should be validated.

    Args:
        shape_name: The shape name to check.

    Returns:
        True if shape should be validated.
    """
    if not shape_name:
        return False

    skip_prefixes = [
        "Title",
        "Subtitle",
        "Content",
        "Picture",
        "Rectangle",
        "Oval",
        "TextBox",
        "Connector",
        "Freeform",
        "Group",
    ]

    for prefix in skip_prefixes:
        if shape_name.startswith(prefix):
            return False

    return "__" in shape_name


def main() -> int:
    """
    Main entry point for command-line usage.

    Returns:
        Exit code (0 for success, 1 for failure).
    """
    if len(sys.argv) < 2:
        print("Usage: python -m backend.tools.validate_template <template.pptx>")
        return 1

    template_path = Path(sys.argv[1])

    try:
        report = validate_template(template_path)
        report.print_report()
        return 0 if report.is_valid() else 1
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
