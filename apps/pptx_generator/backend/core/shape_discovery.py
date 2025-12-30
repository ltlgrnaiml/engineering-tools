"""Shape discovery module per ADR-0019.

Implements the {category}_{identifier}[_{variant}] naming convention for
PowerPoint shape discovery. All mappable shapes MUST have names following
this convention.

Regex pattern: ^(text|chart|table|image|metric|dimension)_([a-zA-Z0-9]+)(?:_([a-zA-Z0-9]+))?$
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide

__version__ = "0.1.0"

# ADR-0019 compliant regex pattern for shape naming
SHAPE_NAME_PATTERN = re.compile(
    r"^(text|chart|table|image|metric|dimension)_([a-zA-Z0-9]+)(?:_([a-zA-Z0-9]+))?$",
    re.IGNORECASE,
)

# Valid categories per ADR-0019
VALID_CATEGORIES = frozenset({"text", "chart", "table", "image", "metric", "dimension"})

# Default PowerPoint shape names to ignore
DEFAULT_SHAPE_NAMES = re.compile(
    r"^(Rectangle|TextBox|Oval|Line|Arrow|Freeform|Picture|Chart|Table|Group|"
    r"Content Placeholder|Title|Subtitle|Footer|Slide Number|Date Placeholder)\s*\d*$",
    re.IGNORECASE,
)


class ShapeNamingError(Exception):
    """Base exception for shape naming errors per ADR-0019."""

    pass


class InvalidCategoryError(ShapeNamingError):
    """Raised when shape name has invalid category."""

    def __init__(self, category: str, shape_name: str) -> None:
        self.category = category
        self.shape_name = shape_name
        valid_list = ", ".join(sorted(VALID_CATEGORIES))
        super().__init__(
            f"Invalid category '{category}' in shape '{shape_name}'. "
            f"Must be one of: {valid_list}"
        )


class InvalidIdentifierError(ShapeNamingError):
    """Raised when shape name has missing or invalid identifier."""

    def __init__(self, shape_name: str, reason: str) -> None:
        self.shape_name = shape_name
        super().__init__(f"Invalid identifier in shape '{shape_name}': {reason}")


class DuplicateShapeNameError(ShapeNamingError):
    """Raised when duplicate shape names found on same slide."""

    def __init__(self, shape_name: str, slide_index: int) -> None:
        self.shape_name = shape_name
        self.slide_index = slide_index
        super().__init__(
            f"Duplicate shape name '{shape_name}' on slide {slide_index + 1}"
        )


@dataclass
class ParsedShapeName:
    """Parsed components of a shape name per ADR-0019 convention.

    Attributes:
        category: Shape category (text, chart, table, image, metric, dimension).
        identifier: Unique identifier within category.
        variant: Optional variant/sub-categorization.
        raw_name: Original unparsed shape name.
    """

    category: str
    identifier: str
    variant: str | None
    raw_name: str

    @property
    def canonical_name(self) -> str:
        """Get canonical (lowercase) shape name."""
        if self.variant:
            return f"{self.category}_{self.identifier}_{self.variant}"
        return f"{self.category}_{self.identifier}"


@dataclass
class DiscoveredShape:
    """Shape discovered from a PowerPoint template.

    Attributes:
        slide_index: Zero-based slide index.
        shape_id: PowerPoint shape ID.
        parsed_name: Parsed shape name components.
        shape_type: PowerPoint shape type (chart, table, text_box, etc.).
        has_text: Whether shape has text frame.
        has_chart: Whether shape has chart.
        has_table: Whether shape has table.
        position: Shape position (left, top, width, height in EMUs).
    """

    slide_index: int
    shape_id: int
    parsed_name: ParsedShapeName
    shape_type: str
    has_text: bool = False
    has_chart: bool = False
    has_table: bool = False
    position: dict[str, int] = field(default_factory=dict)


@dataclass
class ShapeDiscoveryResult:
    """Result of shape discovery from a template per ADR-0019.

    Attributes:
        shapes: List of discovered shapes with valid names.
        errors: List of validation errors (invalid patterns, duplicates).
        warnings: List of validation warnings (orphaned shapes).
        slide_count: Total number of slides in template.
        shapes_by_slide: Shapes grouped by slide index.
    """

    shapes: list[DiscoveredShape]
    errors: list[str]
    warnings: list[str]
    slide_count: int

    @property
    def shapes_by_slide(self) -> dict[int, list[DiscoveredShape]]:
        """Group shapes by slide index."""
        result: dict[int, list[DiscoveredShape]] = {}
        for shape in self.shapes:
            if shape.slide_index not in result:
                result[shape.slide_index] = []
            result[shape.slide_index].append(shape)
        return result

    @property
    def shapes_by_category(self) -> dict[str, list[DiscoveredShape]]:
        """Group shapes by category."""
        result: dict[str, list[DiscoveredShape]] = {}
        for shape in self.shapes:
            cat = shape.parsed_name.category
            if cat not in result:
                result[cat] = []
            result[cat].append(shape)
        return result

    @property
    def is_valid(self) -> bool:
        """Check if discovery completed without errors."""
        return len(self.errors) == 0


def parse_shape_name_adr0018(shape_name: str) -> ParsedShapeName:
    """Parse a shape name per ADR-0019 convention.

    Args:
        shape_name: Shape name to parse.

    Returns:
        ParsedShapeName with extracted components.

    Raises:
        InvalidCategoryError: If category is not valid.
        InvalidIdentifierError: If identifier is missing or invalid.
    """
    if not shape_name or not isinstance(shape_name, str):
        raise InvalidIdentifierError(shape_name or "", "Shape name is empty")

    # Check for underscore separator
    if "_" not in shape_name:
        raise InvalidIdentifierError(
            shape_name,
            "Must use underscore separator: {category}_{identifier}",
        )

    match = SHAPE_NAME_PATTERN.match(shape_name)
    if not match:
        # Try to give helpful error message
        parts = shape_name.split("_")
        if parts[0].lower() not in VALID_CATEGORIES:
            raise InvalidCategoryError(parts[0], shape_name)
        if len(parts) < 2 or not parts[1]:
            raise InvalidIdentifierError(shape_name, "Identifier component is required")
        raise InvalidIdentifierError(
            shape_name,
            "Invalid characters in name. Use only alphanumeric characters and underscores.",
        )

    category = match.group(1).lower()
    identifier = match.group(2).lower()
    variant = match.group(3).lower() if match.group(3) else None

    return ParsedShapeName(
        category=category,
        identifier=identifier,
        variant=variant,
        raw_name=shape_name,
    )


def is_default_shape_name(shape_name: str) -> bool:
    """Check if shape name is a PowerPoint default (should be ignored).

    Per ADR-0019: Reserved PowerPoint default names are automatically ignored.

    Args:
        shape_name: Shape name to check.

    Returns:
        True if this is a default PowerPoint shape name.
    """
    return bool(DEFAULT_SHAPE_NAMES.match(shape_name))


def is_valid_shape_name(shape_name: str) -> bool:
    """Check if shape name follows ADR-0019 convention.

    Args:
        shape_name: Shape name to validate.

    Returns:
        True if name is valid per ADR-0019.
    """
    if not shape_name or is_default_shape_name(shape_name):
        return False
    return bool(SHAPE_NAME_PATTERN.match(shape_name))


def discover_shapes(
    slides: list,
    validate_duplicates: bool = True,
    max_group_depth: int = 10,
) -> ShapeDiscoveryResult:
    """Discover all shapes in a presentation per ADR-0019.

    Implements the discovery algorithm from ADR-0019:
    1. Iterate through slides in order
    2. For each slide, iterate through shapes
    3. Parse and validate shape names
    4. Check for duplicates within slide
    5. Build shape registry

    Args:
        slides: List of pptx Slide objects.
        validate_duplicates: Check for duplicate names per slide.
        max_group_depth: Maximum recursion depth for grouped shapes.

    Returns:
        ShapeDiscoveryResult with discovered shapes, errors, and warnings.
    """
    discovered: list[DiscoveredShape] = []
    errors: list[str] = []
    warnings: list[str] = []

    for slide_idx, slide in enumerate(slides):
        slide_shape_names: set[str] = set()

        def process_shape(shape, depth: int = 0) -> None:
            """Process a single shape, recursing into groups."""
            if depth > max_group_depth:
                warnings.append(
                    f"Slide {slide_idx + 1}: Group recursion depth exceeded {max_group_depth}"
                )
                return

            # Handle grouped shapes per ADR-0019
            if hasattr(shape, "shapes"):
                for child_shape in shape.shapes:
                    process_shape(child_shape, depth + 1)
                return

            shape_name = getattr(shape, "name", None)
            if not shape_name:
                return

            # Skip placeholder shapes per ADR-0019
            if getattr(shape, "is_placeholder", False):
                return

            # Skip default PowerPoint names
            if is_default_shape_name(shape_name):
                return

            # Try to parse shape name
            try:
                parsed = parse_shape_name_adr0018(shape_name)
            except ShapeNamingError as e:
                errors.append(f"Slide {slide_idx + 1}: {e}")
                return

            # Check for duplicates (case-insensitive)
            canonical = parsed.canonical_name
            if validate_duplicates and canonical in slide_shape_names:
                errors.append(
                    f"Slide {slide_idx + 1}: Duplicate shape name '{shape_name}'"
                )
                return
            slide_shape_names.add(canonical)

            # Determine shape type
            shape_type = "unknown"
            has_text = False
            has_chart = False
            has_table = False

            if hasattr(shape, "shape_type"):
                shape_type = str(shape.shape_type).split(".")[-1].lower()

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                has_text = True
            if hasattr(shape, "has_chart") and shape.has_chart:
                has_chart = True
            if hasattr(shape, "has_table") and shape.has_table:
                has_table = True

            # Extract position
            position = {}
            for attr in ("left", "top", "width", "height"):
                if hasattr(shape, attr):
                    position[attr] = getattr(shape, attr)

            discovered.append(
                DiscoveredShape(
                    slide_index=slide_idx,
                    shape_id=getattr(shape, "shape_id", 0),
                    parsed_name=parsed,
                    shape_type=shape_type,
                    has_text=has_text,
                    has_chart=has_chart,
                    has_table=has_table,
                    position=position,
                )
            )

        # Process all shapes on this slide
        for shape in slide.shapes:
            process_shape(shape)

    return ShapeDiscoveryResult(
        shapes=discovered,
        errors=errors,
        warnings=warnings,
        slide_count=len(slides),
    )


def validate_required_shapes(
    discovery_result: ShapeDiscoveryResult,
    required_shapes: list[str],
) -> tuple[list[str], list[str]]:
    """Validate that all required shapes are present.

    Per ADR-0019: Template validation MUST fail if required shapes are missing.

    Args:
        discovery_result: Result from discover_shapes().
        required_shapes: List of required shape names (canonical format).

    Returns:
        Tuple of (missing_shapes, found_shapes).
    """
    discovered_names = {
        shape.parsed_name.canonical_name for shape in discovery_result.shapes
    }

    missing = []
    found = []

    for required in required_shapes:
        # Normalize to lowercase for comparison
        normalized = required.lower()
        if normalized in discovered_names:
            found.append(required)
        else:
            missing.append(required)

    return missing, found
