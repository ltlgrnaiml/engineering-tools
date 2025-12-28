"""API endpoints for slide preview generation.

Per ADR-0031: All errors use ErrorResponse contract via errors.py helper.
"""

import io
import logging
from pathlib import Path
from typing import Any
from uuid import UUID

from fastapi import APIRouter, status
from fastapi.responses import Response
from pptx import Presentation

from apps.pptx_generator.backend.api.errors import (
    raise_not_found,
    raise_validation_error,
    raise_internal_error,
)

from apps.pptx_generator.backend.api.data import data_files_db
from apps.pptx_generator.backend.api.projects import projects_db
from apps.pptx_generator.backend.api.templates import templates_db
from apps.pptx_generator.backend.services.presentation_generator import PresentationGeneratorService

router = APIRouter()
logger = logging.getLogger(__name__)

# In-memory cache for preview images
preview_cache: dict[str, bytes] = {}


@router.get("/{project_id}/preview")
async def get_slide_preview(
    project_id: UUID,
    slide_index: int = 0,
    width: int = 800,
    height: int = 600,
) -> dict[str, Any]:
    """Generate a preview of a slide before full PPTX generation.

    This creates a temporary presentation with rendered content and
    returns information about what would be generated.

    Args:
        project_id: Project ID.
        slide_index: Index of slide to preview (0-based).
        width: Preview image width.
        height: Preview image height.

    Returns:
        Preview information including slide content summary.
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]

    # Check for template
    if not project.template_id or project.template_id not in templates_db:
        raise_validation_error("Template not found. Please upload a template first.", field="template_id")

    template = templates_db[project.template_id]

    # Load template to get slide info
    try:
        prs = Presentation(template.file_path)
    except Exception as e:
        raise_internal_error(f"Failed to load template: {str(e)}", e)

    if slide_index >= len(prs.slides):
        raise_validation_error(
            f"Slide index {slide_index} out of range. Template has {len(prs.slides)} slides.",
            field="slide_index",
        )

    slide = prs.slides[slide_index]

    # Collect shape information for preview
    shapes_info = []
    for shape in slide.shapes:
        shape_info = {
            "name": shape.name,
            "type": shape.shape_type.name if hasattr(shape.shape_type, "name") else str(shape.shape_type),
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        }

        # Add text content if available
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            shape_info["has_text"] = True
            shape_info["text_preview"] = shape.text[:100] if shape.text else ""

        # Check for table
        if hasattr(shape, "has_table") and shape.has_table:
            shape_info["has_table"] = True
            shape_info["table_rows"] = len(shape.table.rows)
            shape_info["table_cols"] = len(shape.table.columns)

        # Check for chart
        if hasattr(shape, "has_chart") and shape.has_chart:
            shape_info["has_chart"] = True

        shapes_info.append(shape_info)

    # Get data preview if available
    data_preview = None
    if project.data_file_id and project.data_file_id in data_files_db:
        data_file = data_files_db[project.data_file_id]
        data_preview = {
            "filename": data_file.filename,
            "row_count": data_file.row_count,
            "columns": data_file.column_names[:10] if data_file.column_names else [],
        }

    return {
        "project_id": str(project_id),
        "slide_index": slide_index,
        "total_slides": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "shapes_count": len(shapes_info),
        "shapes": shapes_info,
        "data_preview": data_preview,
        "preview_available": True,
    }


@router.get("/{project_id}/preview/image")
async def get_slide_preview_image(
    project_id: UUID,
    slide_index: int = 0,
) -> Response:
    """Generate a PNG preview image of a slide.

    Note: This requires additional dependencies (python-pptx-interface or similar)
    that can convert PPTX slides to images. For now, we return a placeholder.

    Args:
        project_id: Project ID.
        slide_index: Index of slide to preview.

    Returns:
        PNG image of the slide preview.
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]

    if not project.template_id or project.template_id not in templates_db:
        raise_validation_error("Template not found.", field="template_id")

    # Check cache first
    cache_key = f"{project_id}_{slide_index}"
    if cache_key in preview_cache:
        return Response(
            content=preview_cache[cache_key],
            media_type="image/png",
        )

    # For now, return slide info as JSON since true image generation
    # requires external tools (LibreOffice, Aspose, etc.)
    # In production, you would use subprocess to call LibreOffice:
    # soffice --headless --convert-to png --outdir /tmp slide.pptx

    raise_validation_error(
        "Image preview requires LibreOffice or similar tool. Use /preview endpoint for shape information."
    )


@router.get("/{project_id}/preview/summary")
async def get_generation_summary(project_id: UUID) -> dict[str, Any]:
    """Get a summary of what will be generated.

    Args:
        project_id: Project ID.

    Returns:
        Summary of generation including slides, shapes, and data.
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]

    summary = {
        "project_id": str(project_id),
        "project_name": project.name,
        "status": project.status.value if hasattr(project.status, "value") else str(project.status),
        "ready_to_generate": False,
        "components": {
            "template": None,
            "data": None,
            "mappings": None,
        },
        "warnings": [],
    }

    # Check template
    if project.template_id and project.template_id in templates_db:
        template = templates_db[project.template_id]
        try:
            prs = Presentation(template.file_path)
            summary["components"]["template"] = {
                "filename": template.filename,
                "slides_count": len(prs.slides),
                "shapes_per_slide": [len(slide.shapes) for slide in prs.slides],
            }
        except Exception as e:
            summary["warnings"].append(f"Could not load template: {str(e)}")
    else:
        summary["warnings"].append("No template uploaded")

    # Check data
    if project.data_file_id and project.data_file_id in data_files_db:
        data_file = data_files_db[project.data_file_id]
        summary["components"]["data"] = {
            "filename": data_file.filename,
            "row_count": data_file.row_count,
            "column_count": len(data_file.column_names) if data_file.column_names else 0,
        }
    else:
        summary["warnings"].append("No data file uploaded")

    # Check mappings
    has_context_mappings = project.context_mapping_id is not None
    has_metrics_mappings = project.metrics_mapping_id is not None
    summary["components"]["mappings"] = {
        "context_configured": has_context_mappings,
        "metrics_configured": has_metrics_mappings,
    }

    if not has_context_mappings:
        summary["warnings"].append("Context mappings not configured")
    if not has_metrics_mappings:
        summary["warnings"].append("Metrics mappings not configured")

    # Determine if ready
    summary["ready_to_generate"] = (
        summary["components"]["template"] is not None
        and summary["components"]["data"] is not None
        and has_context_mappings
        and has_metrics_mappings
        and len(summary["warnings"]) == 0
    )

    return summary
