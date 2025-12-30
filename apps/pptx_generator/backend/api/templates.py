"""Template management endpoints.

Per ADR-0019: Uses ShapeDiscoveryResult contract for shape discovery.
Per ADR-0021: Uses TemplateValidationResult for template validation.
Per ADR-0032: All errors use ErrorResponse contract via errors.py helper.
"""

import time
from datetime import datetime
from pathlib import Path
from uuid import UUID

from fastapi import APIRouter, File, UploadFile, status
from pptx import Presentation

from apps.pptx_generator.backend.api.errors import (
    raise_not_found,
    raise_validation_error,
    raise_internal_error,
)

from shared.contracts.pptx.template import (
    TemplateValidationResult,
    TemplateValidationError,
)

from apps.pptx_generator.backend.api.projects import projects_db
from apps.pptx_generator.backend.core.config import settings
from apps.pptx_generator.backend.core.shape_discovery import (
    discover_shapes,
    ShapeDiscoveryResult,
)
from apps.pptx_generator.backend.models.project import ProjectStatus
from apps.pptx_generator.backend.models.template import ShapeMap, Template
from apps.pptx_generator.backend.services.drm_extractor import DRMExtractorService
from apps.pptx_generator.backend.services.storage import StorageService
from apps.pptx_generator.backend.services.template_parser import TemplateParserService

router = APIRouter()

templates_db: dict[UUID, Template] = {}
shape_maps_db: dict[UUID, ShapeMap] = {}

storage_service = StorageService()
parser_service = TemplateParserService()
drm_extractor = DRMExtractorService()


@router.post("/{project_id}/upload", response_model=Template, status_code=status.HTTP_201_CREATED)
async def upload_template(
    project_id: UUID,
    file: UploadFile = File(...),
) -> Template:
    """
    Upload a PowerPoint template for a project.

    Args:
        project_id: Unique project identifier.
        file: PowerPoint template file (.pptx).

    Returns:
        Template: Uploaded template metadata.

    Raises:
        HTTPException: If project not found (404) or file is invalid (400).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    if not file.filename:
        raise_validation_error("No filename provided", field="file")

    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in settings.ALLOWED_TEMPLATE_EXTENSIONS:
        raise_validation_error(
            f"Invalid file type. Allowed: {settings.ALLOWED_TEMPLATE_EXTENSIONS}",
            field="file",
        )

    template = Template(
        project_id=project_id,
        filename=file.filename,
        file_path="",
        file_size=0,
    )

    try:
        file_path = await storage_service.save_upload(
            file.file,
            template.id,
            file_extension,
        )
        template.file_path = str(file_path)
        template.file_size = file_path.stat().st_size
    except Exception as e:
        raise_internal_error(f"Failed to save template: {str(e)}", e)

    templates_db[template.id] = template

    project = projects_db[project_id]
    project.template_id = template.id
    project.status = ProjectStatus.TEMPLATE_UPLOADED

    from datetime import datetime

    project.updated_at = datetime.utcnow()

    return template


@router.post("/{project_id}/parse", response_model=ShapeMap)
async def parse_template(project_id: UUID) -> ShapeMap:
    """
    Parse a template to extract shape information.

    Args:
        project_id: Unique project identifier.

    Returns:
        ShapeMap: Extracted shape map from the template.

    Raises:
        HTTPException: If project or template not found (404) or parsing fails (500).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.template_id or project.template_id not in templates_db:
        raise_not_found("Template", "No template uploaded for this project")

    template = templates_db[project.template_id]
    template_path = Path(template.file_path)

    try:
        shape_map = await parser_service.parse_template(template_path)
        shape_map.template_id = template.id
    except Exception as e:
        raise_internal_error(f"Failed to parse template: {str(e)}", e)

    shape_maps_db[shape_map.id] = shape_map

    project.shape_map_id = shape_map.id
    project.status = ProjectStatus.TEMPLATE_PARSED

    from datetime import datetime

    project.updated_at = datetime.utcnow()

    return shape_map


@router.post("/{project_id}/parse-v2", response_model=dict)
async def parse_template_v2(project_id: UUID) -> dict:
    """
    Parse a template to extract shape information and DRM (TOM v2).

    This endpoint:
    - Parses all shapes in the template
    - Validates shape names against naming convention
    - Extracts the Derived Requirements Manifest (DRM)
    - Returns validation warnings for problematic shapes

    Args:
        project_id: Unique project identifier.

    Returns:
        Dict containing:
            - shape_map: ShapeMap with all parsed shapes
            - drm: DerivedRequirementsManifest with Four Required Lists
            - warnings: List of validation warnings for malformed shape names

    Raises:
        HTTPException: If project or template not found (404) or parsing fails (500).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.template_id or project.template_id not in templates_db:
        raise_not_found("Template", "No template uploaded for this project")

    template = templates_db[project.template_id]
    template_path = Path(template.file_path)

    try:
        # Use the new parse_template_v2 service method
        shape_map, drm, warnings = await parser_service.parse_template_v2(template_path)
        shape_map.template_id = template.id

        # Store both in databases
        shape_maps_db[shape_map.id] = shape_map

        # Import drm_db from requirements.py (shared database)
        from apps.pptx_generator.backend.api.requirements import drm_db

        drm_db[drm.id] = drm

        # Update project
        project.shape_map_id = shape_map.id
        project.drm_id = drm.id
        project.status = ProjectStatus.DRM_EXTRACTED

        from datetime import datetime

        project.updated_at = datetime.utcnow()

        return {
            "shape_map": shape_map,
            "drm": drm,
            "warnings": warnings,
        }

    except FileNotFoundError as e:
        raise_not_found("Template file", str(e))
    except ValueError as e:
        raise_validation_error(str(e))
    except Exception as e:
        raise_internal_error(f"Failed to parse template: {str(e)}", e)


@router.get("/{project_id}")
async def get_template(project_id: UUID) -> Template:
    """
    Get template information for a project.

    Args:
        project_id: Unique project identifier.

    Returns:
        Template: Template metadata.

    Raises:
        HTTPException: If project or template not found (404).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.template_id or project.template_id not in templates_db:
        raise_not_found("Template", "No template uploaded for this project")

    return templates_db[project.template_id]


@router.get("/{project_id}/shape-map", response_model=ShapeMap)
async def get_shape_map(project_id: UUID) -> ShapeMap:
    """
    Get shape map for a project.

    Args:
        project_id: Unique project identifier.

    Returns:
        ShapeMap: Shape map with all template shapes.

    Raises:
        HTTPException: If project or shape map not found (404).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.shape_map_id or project.shape_map_id not in shape_maps_db:
        raise_not_found("ShapeMap", "No shape map available. Please parse the template first.")

    return shape_maps_db[project.shape_map_id]


@router.post("/{project_id}/discover-shapes")
async def discover_template_shapes(project_id: UUID) -> dict:
    """Discover shapes in template per ADR-0019.

    Uses the ShapeDiscoveryResult contract for comprehensive shape analysis.

    Args:
        project_id: Unique project identifier.

    Returns:
        Dict containing ShapeDiscoveryResult data:
            - shapes: List of discovered shapes with parsed names.
            - errors: Validation errors (invalid patterns, duplicates).
            - warnings: Validation warnings.
            - slide_count: Total slides in template.
            - is_valid: True if no errors.

    Raises:
        HTTPException: If project or template not found (404).
    """
    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.template_id or project.template_id not in templates_db:
        raise_not_found("Template", "No template uploaded for this project")

    template = templates_db[project.template_id]
    template_path = Path(template.file_path)

    try:
        prs = Presentation(str(template_path))
        discovery_result = discover_shapes(list(prs.slides))

        return {
            "shapes": [
                {
                    "slide_index": s.slide_index,
                    "shape_id": s.shape_id,
                    "category": s.parsed_name.category,
                    "identifier": s.parsed_name.identifier,
                    "variant": s.parsed_name.variant,
                    "raw_name": s.parsed_name.raw_name,
                    "canonical_name": s.parsed_name.canonical_name,
                    "shape_type": s.shape_type,
                    "has_text": s.has_text,
                    "has_chart": s.has_chart,
                    "has_table": s.has_table,
                }
                for s in discovery_result.shapes
            ],
            "errors": discovery_result.errors,
            "warnings": discovery_result.warnings,
            "slide_count": discovery_result.slide_count,
            "is_valid": discovery_result.is_valid,
            "shapes_by_category": {
                cat: len(shapes)
                for cat, shapes in discovery_result.shapes_by_category.items()
            },
        }

    except Exception as e:
        raise_internal_error(f"Failed to discover shapes: {str(e)}", e)


@router.post("/{project_id}/validate")
async def validate_template(project_id: UUID) -> dict:
    """Validate a template per ADR-0019/0020.

    Uses the TemplateValidationResult contract for comprehensive validation.

    Args:
        project_id: Unique project identifier.

    Returns:
        TemplateValidationResult with validation status and errors.

    Raises:
        HTTPException: If project or template not found (404).
    """
    start_time = time.time()

    if project_id not in projects_db:
        raise_not_found("Project", str(project_id))

    project = projects_db[project_id]
    if not project.template_id or project.template_id not in templates_db:
        raise_not_found("Template", "No template uploaded for this project")

    template = templates_db[project.template_id]
    template_path = Path(template.file_path)

    errors: list[TemplateValidationError] = []
    warnings: list[TemplateValidationError] = []
    layouts_found = 0
    placeholders_found = 0
    shapes_found = 0

    try:
        prs = Presentation(str(template_path))

        # Count layouts
        layouts_found = len(prs.slide_layouts)

        # Discover and validate shapes
        discovery_result = discover_shapes(list(prs.slides))
        shapes_found = len(discovery_result.shapes)

        # Convert discovery errors to validation errors
        for err in discovery_result.errors:
            errors.append(TemplateValidationError(
                error_code="SHAPE_DISCOVERY_ERROR",
                message=err,
                severity="error",
            ))

        for warn in discovery_result.warnings:
            warnings.append(TemplateValidationError(
                error_code="SHAPE_DISCOVERY_WARNING",
                message=warn,
                severity="warning",
            ))

        # Count placeholders
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholders_found += 1

        # Check for required shapes
        if shapes_found == 0:
            warnings.append(TemplateValidationError(
                error_code="NO_MAPPABLE_SHAPES",
                message="No shapes with ADR-0019 compliant names found",
                severity="warning",
            ))

    except Exception as e:
        errors.append(TemplateValidationError(
            error_code="TEMPLATE_PARSE_ERROR",
            message=f"Failed to parse template: {str(e)}",
            severity="error",
        ))

    validation_duration_ms = (time.time() - start_time) * 1000

    result = TemplateValidationResult(
        template_id=str(template.id),
        valid=len(errors) == 0,
        errors=errors,
        warnings=warnings,
        layouts_found=layouts_found,
        placeholders_found=placeholders_found,
        shapes_found=shapes_found,
        validation_duration_ms=validation_duration_ms,
    )

    return result.model_dump()
