"""Presentation generator service for creating PowerPoint presentations."""

import logging
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation

from apps.pptx_generator.backend.core.shape_name_parser import parse_shape_name
from apps.pptx_generator.backend.renderers.base import RenderContext
from apps.pptx_generator.backend.renderers.factory import RendererFactory
from apps.pptx_generator.backend.renderers.inert_renderer import ImageRenderer, InertRenderer
from apps.pptx_generator.backend.renderers.plot_renderer import PlotRenderer
from apps.pptx_generator.backend.renderers.table_renderer import TableRenderer
from apps.pptx_generator.backend.renderers.text_renderer import KPIRenderer, TextRenderer

logger = logging.getLogger(__name__)


class PresentationGeneratorService:
    """
    Service for generating PowerPoint presentations from templates and data.

    Handles the core logic of populating templates with user data using
    a modular renderer system.
    """

    def __init__(self):
        """Initialize the presentation generator with renderer factory."""
        self.renderer_factory = RendererFactory()
        self._register_renderers()
        self.logger = logging.getLogger(__name__)

    def _register_renderers(self) -> None:
        """Register all available renderers."""
        self.renderer_factory.register_renderer(TextRenderer())
        self.renderer_factory.register_renderer(KPIRenderer())
        self.renderer_factory.register_renderer(TableRenderer())
        self.renderer_factory.register_renderer(PlotRenderer())
        self.renderer_factory.register_renderer(ImageRenderer())
        self.renderer_factory.register_renderer(InertRenderer())

    async def generate_presentation(
        self,
        template_path: Path,
        data_records: list[dict[str, Any]],
        output_path: Path,
    ) -> Path:
        """
        Generate a PowerPoint presentation from a template and data.

        Args:
            template_path: Path to the PowerPoint template file.
            data_records: List of data records to populate in the presentation.
            output_path: Path where the generated presentation should be saved.

        Returns:
            Path: Path to the generated presentation file.

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template or data is invalid.
            IOError: If presentation cannot be saved.
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

        # Convert data records to DataFrame
        data_df = pd.DataFrame(data_records) if data_records else pd.DataFrame()

        self.logger.info(f"Processing presentation with {len(data_df)} data records")

        # DEBUG: Log DataFrame structure
        if not data_df.empty:
            self.logger.info(f"DataFrame columns: {list(data_df.columns)}")
            self.logger.info(f"DataFrame shape: {data_df.shape}")
            self.logger.info(f"DataFrame head:\n{data_df.head()}")
            # Also print to ensure visibility
            print(f"[PRESENTATION_GEN] DataFrame columns: {list(data_df.columns)}")
            print(f"[PRESENTATION_GEN] DataFrame shape: {data_df.shape}")
            print(f"[PRESENTATION_GEN] First row: {data_df.iloc[0].to_dict()}")

        # Populate slides with renderer system
        await self._populate_slides_with_renderers(prs, data_df, output_path.parent)

        try:
            prs.save(str(output_path))
        except Exception as e:
            raise OSError(f"Failed to save presentation: {str(e)}") from e

        self.logger.info(f"Presentation saved to {output_path}")
        return output_path

    async def _populate_slides_with_renderers(
        self,
        presentation: Presentation,
        data: pd.DataFrame,
        output_dir: Path,
    ) -> None:
        """Populate all slides using renderer system.

        Args:
            presentation: PowerPoint presentation object.
            data: Data as DataFrame.
            output_dir: Directory for temporary files.
        """
        total_shapes = 0
        rendered_shapes = 0

        for slide_idx, slide in enumerate(presentation.slides):
            self.logger.debug(f"Processing slide {slide_idx + 1}/{len(presentation.slides)}")

            # CRITICAL: Snapshot shapes before iteration to prevent infinite loop
            # Renderers may add/remove shapes, so we must process original shapes only
            shapes_snapshot = list(slide.shapes)
            self.logger.debug(f"Slide {slide_idx + 1} has {len(shapes_snapshot)} shapes to process")

            for shape in shapes_snapshot:
                total_shapes += 1

                # Try to parse shape name
                try:
                    parsed_name = parse_shape_name(shape.name)
                except Exception as e:
                    self.logger.debug(f"Skipping unparseable shape '{shape.name}': {e}")
                    continue

                # Get appropriate renderer
                renderer = self.renderer_factory.get_renderer(parsed_name)

                if not renderer:
                    self.logger.debug(f"No renderer for shape '{shape.name}'")
                    continue

                # Create render context
                context = RenderContext(
                    shape=shape,
                    parsed_name=parsed_name,
                    data=data,
                    output_dir=output_dir,
                )

                # Render
                try:
                    await renderer.render(context)
                    rendered_shapes += 1
                except Exception as e:
                    self.logger.error(f"Failed to render shape '{shape.name}': {e}", exc_info=True)

        self.logger.info(
            f"Rendered {rendered_shapes}/{total_shapes} shapes across "
            f"{len(presentation.slides)} slides"
        )

    def _populate_slides(self, presentation: Presentation, data: dict[str, Any]) -> None:
        """
        Populate all slides in a presentation with data.

        DEPRECATED: Use _populate_slides_with_renderers instead.

        Args:
            presentation: PowerPoint presentation object.
            data: Dictionary mapping shape names to values.
        """
        for slide in presentation.slides:
            self._populate_slide(slide, data)

    def _populate_slide(self, slide: Any, data: dict[str, Any]) -> None:
        """
        Populate a single slide with data.

        Args:
            slide: PowerPoint slide object.
            data: Dictionary mapping shape names to values.
        """
        for shape in slide.shapes:
            shape_name = shape.name if hasattr(shape, "name") else None

            if shape_name and shape_name in data:
                value = data[shape_name]
                self._set_shape_value(shape, value)

    def _set_shape_value(self, shape: Any, value: Any) -> None:
        """
        Set the value of a shape based on its type.

        Args:
            shape: PowerPoint shape object.
            value: Value to set in the shape.
        """
        try:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if shape.text_frame.text or not shape.text_frame.paragraphs:
                    shape.text = str(value) if value is not None else ""
                else:
                    shape.text_frame.paragraphs[0].text = str(value) if value is not None else ""

            elif (
                hasattr(shape, "has_table")
                and shape.has_table
                or hasattr(shape, "has_chart")
                and shape.has_chart
            ):
                pass

        except Exception:
            pass

    async def duplicate_slide(
        self,
        presentation: Presentation,
        slide_index: int,
    ) -> None:
        """
        Duplicate a slide in a presentation.

        Args:
            presentation: PowerPoint presentation object.
            slide_index: Index of the slide to duplicate.

        Raises:
            IndexError: If slide_index is out of range.
        """
        if slide_index < 0 or slide_index >= len(presentation.slides):
            raise IndexError(f"Slide index {slide_index} out of range")

        source_slide = presentation.slides[slide_index]
        blank_slide_layout = presentation.slide_layouts[6]
        new_slide = presentation.slides.add_slide(blank_slide_layout)

        for shape in source_slide.shapes:
            el = shape.element
            newel = el.__class__(el)
            new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

    async def generate_multi_slide_presentation(
        self,
        template_path: Path,
        data_records: list[dict[str, Any]],
        output_path: Path,
        slide_template_index: int = 0,
    ) -> Path:
        """
        Generate a multi-slide presentation with one slide per data record.

        Args:
            template_path: Path to the PowerPoint template file.
            data_records: List of data records, one per slide.
            output_path: Path where the generated presentation should be saved.
            slide_template_index: Index of the slide to use as template (default: 0).

        Returns:
            Path: Path to the generated presentation file.

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template or data is invalid.
            IOError: If presentation cannot be saved.
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

        if slide_template_index >= len(prs.slides):
            raise ValueError(
                f"Slide template index {slide_template_index} out of range. "
                f"Template has {len(prs.slides)} slides."
            )

        template_slide = prs.slides[slide_template_index]

        for idx, record in enumerate(data_records):
            if idx == 0:
                self._populate_slide(template_slide, record)
            else:
                await self.duplicate_slide(prs, slide_template_index)
                new_slide = prs.slides[-1]
                self._populate_slide(new_slide, record)

        try:
            prs.save(str(output_path))
        except Exception as e:
            raise OSError(f"Failed to save presentation: {str(e)}") from e

        return output_path
