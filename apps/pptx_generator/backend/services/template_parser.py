"""Template parser service for extracting shape information from PowerPoint templates."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from apps.pptx_generator.backend.core.shape_discovery import (
    ShapeDiscoveryResult,
    discover_shapes,
)
from apps.pptx_generator.backend.core.shape_name_parser import ValidationError, parse_shape_name
from apps.pptx_generator.backend.models.drm import DerivedRequirementsManifest
from apps.pptx_generator.backend.models.template import ShapeInfo, ShapeMap
from apps.pptx_generator.backend.services.drm_extractor import DRMExtractorService


class TemplateParserService:
    """
    Service for parsing PowerPoint templates and extracting shape information.

    Analyzes template files to create shape maps that can be used for data mapping.
    """

    @staticmethod
    def _get_shape_type_name(shape_type: MSO_SHAPE_TYPE) -> str:
        """
        Convert MSO_SHAPE_TYPE enum to human-readable string.

        Args:
            shape_type: PowerPoint shape type enum.

        Returns:
            str: Human-readable shape type name.
        """
        type_mapping = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
            MSO_SHAPE_TYPE.CALLOUT: "callout",
            MSO_SHAPE_TYPE.CANVAS: "canvas",
            MSO_SHAPE_TYPE.CHART: "chart",
            MSO_SHAPE_TYPE.COMMENT: "comment",
            MSO_SHAPE_TYPE.FREEFORM: "freeform",
            MSO_SHAPE_TYPE.GROUP: "group",
            MSO_SHAPE_TYPE.LINE: "line",
            MSO_SHAPE_TYPE.MEDIA: "media",
            MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: "ole_control",
            MSO_SHAPE_TYPE.PICTURE: "picture",
            MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
            MSO_SHAPE_TYPE.SCRIPT_ANCHOR: "script_anchor",
            MSO_SHAPE_TYPE.TABLE: "table",
            MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
        }
        return type_mapping.get(shape_type, "unknown")

    async def parse_template(self, template_path: Path) -> ShapeMap:
        """
        Parse a PowerPoint template and extract all shape information.

        Args:
            template_path: Path to the PowerPoint template file.

        Returns:
            ShapeMap: Complete mapping of all shapes in the template.

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template file is invalid or corrupted.
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

        shapes: list[ShapeInfo] = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                shape_info = self._extract_shape_info(shape, slide_idx)
                if shape_info:
                    shapes.append(shape_info)

        from uuid import uuid4

        shape_map = ShapeMap(
            template_id=uuid4(),
            shapes=shapes,
            slide_count=len(prs.slides),
        )

        return shape_map

    def _extract_shape_info(self, shape, slide_index: int) -> ShapeInfo | None:
        """
        Extract information from a single shape.

        Args:
            shape: PowerPoint shape object.
            slide_index: Zero-based index of the slide containing the shape.

        Returns:
            ShapeInfo | None: Shape information or None if shape should be skipped.
        """
        try:
            shape_name = shape.name if hasattr(shape, "name") else f"Shape_{shape.shape_id}"

            has_text = False
            has_table = False
            has_chart = False

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                has_text = True

            if hasattr(shape, "has_table") and shape.has_table:
                has_table = True

            if hasattr(shape, "has_chart") and shape.has_chart:
                has_chart = True

            shape_type_name = self._get_shape_type_name(shape.shape_type)

            position = {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            return ShapeInfo(
                shape_id=shape.shape_id,
                name=shape_name,
                shape_type=shape_type_name,
                slide_index=slide_index,
                position=position,
                has_text=has_text,
                has_table=has_table,
                has_chart=has_chart,
            )
        except Exception:
            return None

    async def get_shape_names(self, template_path: Path) -> list[str]:
        """
        Get a list of all shape names in a template.

        Args:
            template_path: Path to the PowerPoint template file.

        Returns:
            List[str]: List of unique shape names.

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template file is invalid.
        """
        shape_map = await self.parse_template(template_path)
        return [shape.name for shape in shape_map.shapes]

    async def discover_shapes_adr0018(
        self,
        template_path: Path,
        required_shapes: list[str] | None = None,
    ) -> ShapeDiscoveryResult:
        """Discover shapes using ADR-0019 naming convention.

        Per ADR-0019: Uses {category}_{identifier}[_{variant}] convention.
        Regex: ^(text|chart|table|image|metric|dimension)_([a-zA-Z0-9]+)(?:_([a-zA-Z0-9]+))?$

        Args:
            template_path: Path to the PowerPoint template file.
            required_shapes: Optional list of required shape names to validate.

        Returns:
            ShapeDiscoveryResult with discovered shapes, errors, and warnings.

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template file is invalid.
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

        # Use the ADR-0019 compliant shape discovery
        result = discover_shapes(list(prs.slides))

        # Validate required shapes if provided
        if required_shapes:
            from apps.pptx_generator.backend.core.shape_discovery import validate_required_shapes
            missing, _ = validate_required_shapes(result, required_shapes)
            if missing:
                result.errors.extend([
                    f"Missing required shape: {name}" for name in missing
                ])

        return result

    async def parse_template_v2(
        self, template_path: Path
    ) -> tuple[ShapeMap, DerivedRequirementsManifest, list[str]]:
        """
        Parse a PowerPoint template and extract shape map AND DRM (TOM v2).

        This method extends the basic template parsing to also:
        - Parse shape names using the shape naming convention
        - Extract Derived Requirements Manifest (DRM)
        - Collect validation warnings for problematic shape names

        Args:
            template_path: Path to the PowerPoint template file.

        Returns:
            Tuple containing:
                - ShapeMap: Complete mapping of all shapes
                - DerivedRequirementsManifest: Extracted requirements
                - List[str]: Validation warnings for malformed shape names

        Raises:
            FileNotFoundError: If template file doesn't exist.
            ValueError: If template file is invalid or corrupted.
        """
        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint template: {str(e)}") from e

        shapes: list[ShapeInfo] = []
        warnings: list[str] = []

        # Parse all shapes with validation
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                shape_info = self._extract_shape_info_v2(shape, slide_idx, warnings)
                if shape_info:
                    shapes.append(shape_info)

        from uuid import uuid4

        shape_map = ShapeMap(
            template_id=uuid4(),
            shapes=shapes,
            slide_count=len(prs.slides),
        )

        # Extract DRM using the DRM extractor service
        drm_extractor = DRMExtractorService()
        drm = drm_extractor.extract_drm(shape_map)

        return shape_map, drm, warnings

    def _extract_shape_info_v2(
        self, shape, slide_index: int, warnings: list[str]
    ) -> ShapeInfo | None:
        """
        Extract information from a single shape with shape name parsing (TOM v2).

        Args:
            shape: PowerPoint shape object.
            slide_index: Zero-based index of the slide containing the shape.
            warnings: List to append validation warnings to.

        Returns:
            ShapeInfo | None: Shape information or None if shape should be skipped.
        """
        try:
            shape_name = shape.name if hasattr(shape, "name") else f"Shape_{shape.shape_id}"

            # Try to parse the shape name
            parsed_name = None
            try:
                parsed_name = parse_shape_name(shape_name)
            except ValidationError as e:
                warnings.append(f"Slide {slide_index + 1}, Shape '{shape_name}': {str(e)}")
            except Exception as e:
                warnings.append(
                    f"Slide {slide_index + 1}, Shape '{shape_name}': "
                    f"Unexpected parsing error: {str(e)}"
                )

            has_text = False
            has_table = False
            has_chart = False

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                has_text = True

            if hasattr(shape, "has_table") and shape.has_table:
                has_table = True

            if hasattr(shape, "has_chart") and shape.has_chart:
                has_chart = True

            shape_type_name = self._get_shape_type_name(shape.shape_type)

            position = {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            return ShapeInfo(
                shape_id=shape.shape_id,
                name=shape_name,
                shape_type=shape_type_name,
                slide_index=slide_index,
                position=position,
                has_text=has_text,
                has_table=has_table,
                has_chart=has_chart,
                parsed_name=parsed_name,
            )
        except Exception:
            return None
